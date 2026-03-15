import hashlib
import io
from typing import BinaryIO

from docx import Document
import olefile
from pptx import Presentation
from pypdf import PdfReader


def get_file_extension(filename: str) -> str:
    index = filename.rfind(".")
    if index == -1:
        raise ValueError(f"Can not find extension in filename '{filename}'")
    return filename[index::]


def sha256_checksum(file: BinaryIO) -> str:
    sha256 = hashlib.sha256()
    while chunk := file.read(1024 * 1024):
        sha256.update(chunk)
    file.seek(0)
    return sha256.hexdigest()


def md5_checksum(file: BinaryIO) -> str:
    md5 = hashlib.md5()
    while chunk := file.read(1024 * 1024):
        md5.update(chunk)
    file.seek(0)
    return md5.hexdigest()


def get_page_count(binary_stream, extension):
    extension = extension.lower().strip(".")

    if extension == "pdf":
        reader = PdfReader(binary_stream)
        return len(reader.pages)

    elif extension == "docx":
        doc = Document(binary_stream)

        return doc.core_properties.pages if doc.core_properties.pages else 0

    elif extension == "pptx":
        prs = Presentation(binary_stream)
        return len(prs.slides)

    elif extension in ["doc", "ppt"]:
        ole = olefile.OleFileIO(binary_stream)
        meta = ole.get_metadata()
        return meta.num_pages if extension == "doc" else meta.slides

    return 0
